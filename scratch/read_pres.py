from pptx import Presentation

file_path = "E:\\Dip\\gsdfg\\New\\DroneLocalization\\Презентація захист.pptx"
try:
    prs = Presentation(file_path)
    print(f"Total slides: {len(prs.slides)}")
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                print(shape.text)
        print()
except Exception as e:
    print(f"Error reading presentation: {e}")
